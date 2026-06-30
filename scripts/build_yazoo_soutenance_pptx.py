from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\seef7\OneDrive\Desktop\YaZoo")
OUT_DIR = ROOT / "docs" / "soutenance"
OUT = OUT_DIR / "YaZoo_Soutenance_PFE_DEVFS201.pptx"

LOGO_YAZOO = Path(r"C:\Users\seef7\Downloads\YaZoo_Logo_Cat-Bird_Hybrid_with_Paw_Print__1_-removebg-preview.png")
LOGO_OFPPT = Path(r"C:\Users\seef7\Downloads\unnamed.jpg")
IMG_DIR = Path(r"C:\Users\seef7\OneDrive\Pictures\siteYazoo")
SHOTS = [
    IMG_DIR / "Capture d’écran 2026-05-19 211856.png",
    IMG_DIR / "Capture d’écran 2026-05-19 211950.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212057.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212109.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212258.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212342.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212401.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212527.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212603.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212644.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212707.png",
    IMG_DIR / "Capture d’écran 2026-05-19 212758.png",
]


PURPLE = RGBColor(123, 58, 237)
VIOLET = RGBColor(168, 85, 247)
LILAC = RGBColor(244, 236, 255)
INK = RGBColor(28, 25, 23)
MUTED = RGBColor(92, 82, 101)
LINE = RGBColor(222, 210, 245)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(22, 163, 74)
BLUE = RGBColor(37, 99, 235)
ORANGE = RGBColor(234, 88, 12)


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("fffaff")
    shape.line.fill.background()
    # Soft violet bands, simple and editable.
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = LILAC
    band.line.fill.background()
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.28), Inches(13.333), Inches(0.22))
    foot.fill.solid()
    foot.fill.fore_color.rgb = LILAC
    foot.line.fill.background()


def set_text(tf, text, size=16, bold=False, color=INK, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align


def text_box(slide, x, y, w, h, text, size=16, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = Inches(0.02)
    box.text_frame.margin_right = Inches(0.02)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_text(box.text_frame, text, size, bold, color, align)
    return box


def title(slide, kicker, claim):
    text_box(slide, 0.62, 0.42, 3.8, 0.28, kicker.upper(), 12, True, PURPLE)
    return text_box(slide, 0.62, 0.75, 9.8, 0.55, claim, 27, True, INK)


def footer(slide, n):
    text_box(slide, 0.62, 7.08, 5.6, 0.2, "YaZoo - Projet de fin de formation | DEVFS 201", 9, False, MUTED)
    text_box(slide, 12.4, 7.08, 0.35, 0.2, str(n), 9, True, PURPLE, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def pill(slide, x, y, w, h, label, fill=PURPLE, color=WHITE):
    shp = card(slide, x, y, w, h, fill, fill)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shp.text_frame, label, 13, True, color, PP_ALIGN.CENTER)
    return shp


def bullet_list(slide, x, y, w, h, items, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)
        p.text = "• " + item
    return box


def picture(slide, path: Path, x, y, w, h):
    if path.exists():
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    return card(slide, x, y, w, h, LILAC)


def line(slide, x1, y1, x2, y2, color=PURPLE, width=1.5, arrow=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def node(slide, x, y, w, h, label, fill=WHITE, border=PURPLE, size=14, bold=True):
    shp = card(slide, x, y, w, h, fill, border)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shp.text_frame, label, size, bold, INK, PP_ALIGN.CENTER)
    return shp


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    card(s, 0.55, 0.45, 12.25, 6.55, WHITE, LINE)
    picture(s, LOGO_OFPPT, 0.83, 0.78, 1.0, 1.0)
    picture(s, LOGO_YAZOO, 11.15, 0.72, 1.05, 1.05)
    text_box(s, 2.0, 0.85, 4.4, 0.45, "OFPPT - Développement Digital", 16, True, MUTED)
    text_box(s, 2.0, 1.22, 4.4, 0.35, "Option Full-stack | Groupe DEVFS 201", 16, False, MUTED)
    text_box(s, 1.05, 2.05, 6.2, 0.75, "YaZoo", 28, True, PURPLE)
    text_box(s, 1.05, 2.78, 6.9, 1.05, "Plateforme sociale et marketplace animalière", 26, True, INK)
    text_box(s, 1.05, 3.9, 5.7, 0.7, "Une application web moderne pour partager, adopter, vendre, réserver et échanger autour des animaux.", 17, False, MUTED)
    pill(s, 1.05, 4.85, 2.7, 0.45, "Soutenance PFE")
    text_box(s, 1.05, 5.55, 4.8, 0.35, "Présenté par :", 16, True, INK)
    text_box(s, 1.05, 5.92, 5.9, 0.72, "Youssef Boughioul | Salma Kabran | Chams Doha Amine", 16, False, INK)
    text_box(s, 7.3, 5.55, 3.8, 0.35, "Encadrante :", 16, True, INK)
    text_box(s, 7.3, 5.92, 3.8, 0.35, "Imane LAASSAR", 16, False, INK)
    picture(s, SHOTS[0], 7.15, 2.05, 4.9, 2.75)
    footer(s, 1)


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Plan", "La soutenance suit le chemin du besoin vers la solution déployée.")
    items = [
        ("01", "Contexte et problématique"),
        ("02", "Objectifs et périmètre fonctionnel"),
        ("03", "Analyse UML et conception"),
        ("04", "Réalisation technique Laravel / React"),
        ("05", "Docker, CI/CD et déploiement Azure"),
        ("06", "Tests, bilan et perspectives"),
    ]
    for i, (num, label) in enumerate(items):
        x = 0.95 + (i % 2) * 5.9
        y = 1.75 + (i // 2) * 1.35
        card(s, x, y, 5.35, 0.92, WHITE, LINE)
        pill(s, x + 0.18, y + 0.23, 0.7, 0.38, num, LILAC, PURPLE)
        text_box(s, x + 1.05, y + 0.22, 3.9, 0.4, label, 17, True, INK)
    footer(s, 2)


def slide_context(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Contexte", "Les passionnés d'animaux ont besoin d'un espace unique et fiable.")
    bullet_list(s, 0.8, 1.55, 5.4, 3.0, [
        "Les publications, annonces, messages et réservations sont souvent dispersés.",
        "Les adoptions et ventes nécessitent plus de visibilité et de confiance.",
        "Les communautés locales manquent d'outils simples pour organiser l'entraide.",
        "YaZoo centralise ces usages dans une plateforme sociale 100% orientée animaux.",
    ])
    picture(s, SHOTS[1], 6.75, 1.45, 5.6, 3.15)
    card(s, 6.75, 4.92, 5.6, 0.75, LILAC, LINE)
    text_box(s, 7.02, 5.08, 5.0, 0.35, "Vision: créer une communauté animalière marocaine, utile et rassurante.", 17, True, PURPLE)
    footer(s, 3)


def slide_objectives(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Objectifs", "YaZoo réunit réseau social, marketplace et suivi des réservations.")
    objs = [
        ("Social", "Publier posts, photos, vidéos, stories, likes et commentaires."),
        ("Marketplace", "Créer des annonces animaux et produits avec images et filtres."),
        ("Communautés", "Rechercher, créer, rejoindre et modérer des groupes."),
        ("Messagerie", "Discuter entre utilisateurs autour des annonces ou adoptions."),
        ("Réservations", "Approuver, rejeter, livrer, terminer et générer une facture."),
        ("Administration", "Modération globale et tableau des commandes."),
    ]
    for i, (h, b) in enumerate(objs):
        x = 0.75 + (i % 3) * 4.15
        y = 1.55 + (i // 3) * 1.9
        card(s, x, y, 3.75, 1.35, WHITE, LINE)
        text_box(s, x + 0.22, y + 0.18, 2.9, 0.28, h, 18, True, PURPLE)
        text_box(s, x + 0.22, y + 0.55, 3.25, 0.55, b, 16, False, MUTED)
    footer(s, 4)


def slide_stack(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Technologies", "Le projet combine une stack web moderne et une architecture cloud.")
    cols = [
        ("Frontend", ["React 19", "Vite 8", "React Router", "Axios", "Tailwind CSS", "Vitest / ESLint"], PURPLE),
        ("Backend", ["Laravel 12", "PHP 8.4", "Sanctum", "MySQL", "Redis queues/cache", "PHPUnit"], BLUE),
        ("DevOps", ["Docker", "Nginx", "GitHub Actions", "Docker Hub", "SonarQube", "Azure App Service"], GREEN),
    ]
    for i, (h, items, color) in enumerate(cols):
        x = 0.8 + i * 4.15
        card(s, x, 1.55, 3.75, 4.45, WHITE, LINE)
        pill(s, x + 0.25, 1.85, 1.55, 0.42, h, color)
        bullet_list(s, x + 0.35, 2.55, 3.0, 2.85, items, 16)
    footer(s, 5)


def slide_use_case(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "UML - Cas d'utilisation", "Trois acteurs couvrent les parcours publics, membres et administration.")
    node(s, 0.65, 2.2, 1.55, 0.55, "Visiteur", LILAC)
    node(s, 0.65, 3.35, 1.55, 0.55, "Utilisateur", LILAC)
    node(s, 0.65, 4.5, 1.55, 0.55, "Admin", LILAC)
    use_cases = [
        ("Consulter accueil", 3.0, 1.35),
        ("S'inscrire / se connecter", 3.0, 2.25),
        ("Publier post / story", 5.65, 1.35),
        ("Gérer marketplace", 5.65, 2.25),
        ("Réserver / commander", 5.65, 3.15),
        ("Créer communauté", 8.3, 1.35),
        ("Messagerie", 8.3, 2.25),
        ("Notifications", 8.3, 3.15),
        ("Modérer contenus", 5.65, 4.45),
        ("Tableau commandes", 8.3, 4.45),
    ]
    for label, x, y in use_cases:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(2.05), Inches(0.58))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = PURPLE
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(shp.text_frame, label, 13, False, INK, PP_ALIGN.CENTER)
    for yy, targets in [(2.47, [(4.0, 1.64), (4.0, 2.54)]), (3.62, [(6.65, 1.64), (6.65, 2.54), (6.65, 3.44), (9.3, 1.64), (9.3, 2.54), (9.3, 3.44)]), (4.77, [(6.65, 4.74), (9.3, 4.74)])]:
        for tx, ty in targets:
            line(s, 2.2, yy, tx, ty, LINE, 1)
    footer(s, 6)


def slide_class(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "UML - Diagramme de classes", "Le modèle métier s'articule autour de l'utilisateur et de ses interactions.")
    classes = {
        "User": (0.7, 1.55, ["id", "name", "email", "is_admin"]),
        "Post": (3.2, 1.25, ["caption", "media_path", "location"]),
        "Comment": (5.7, 1.25, ["content", "user_id", "post_id"]),
        "Like": (8.2, 1.25, ["user_id", "post_id"]),
        "Story": (10.4, 1.25, ["media", "expires_at"]),
        "Animal": (2.0, 4.0, ["name", "category", "price", "status"]),
        "Product": (4.55, 4.0, ["name", "stock", "price"]),
        "Reservation": (7.1, 4.0, ["status", "quantity", "invoice"]),
        "Community": (9.8, 4.0, ["name", "private", "owner_id"]),
    }
    for c, (x, y, attrs) in classes.items():
        card(s, x, y, 2.05, 1.08, WHITE, LINE, MSO_SHAPE.RECTANGLE)
        text_box(s, x + 0.08, y + 0.08, 1.8, 0.22, c, 14, True, PURPLE, PP_ALIGN.CENTER)
        text_box(s, x + 0.12, y + 0.38, 1.78, 0.56, "\n".join(attrs), 11, False, MUTED)
    for a, b in [
        ((1.75, 2.1), (4.2, 1.8)), ((1.75, 2.1), (6.7, 1.8)),
        ((1.75, 2.1), (9.2, 1.8)), ((1.75, 2.1), (11.4, 1.8)),
        ((1.75, 2.65), (3.0, 4.55)), ((1.75, 2.65), (5.55, 4.55)),
        ((3.05, 4.55), (8.1, 4.55)), ((5.6, 4.55), (8.1, 4.55)),
        ((1.75, 2.65), (10.8, 4.55)),
    ]:
        line(s, a[0], a[1], b[0], b[1], LINE, 1)
    footer(s, 7)


def slide_sequence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "UML - Séquence", "La connexion sécurisée utilise React, Axios, Laravel Sanctum et la base MySQL.")
    actors = [("Utilisateur", 0.8), ("React UI", 3.1), ("Axios API", 5.1), ("Laravel Auth", 7.2), ("MySQL", 10.0)]
    for label, x in actors:
        node(s, x, 1.45, 1.5, 0.45, label, LILAC)
        line(s, x + 0.75, 1.95, x + 0.75, 5.95, LINE, 1)
    steps = [
        (1.55, 3.85, "email + mot de passe"),
        (3.85, 5.85, "POST /auth/login"),
        (5.85, 7.95, "LoginRequest + AuthService"),
        (7.95, 10.75, "vérifier utilisateur"),
        (10.75, 7.95, "résultat"),
        (7.95, 5.85, "token Sanctum cookie"),
        (5.85, 3.85, "session authentifiée"),
    ]
    y = 2.35
    for x1, x2, label in steps:
        line(s, x1, y, x2, y, PURPLE if x2 > x1 else GREEN, 1.4, True)
        text_box(s, min(x1, x2) + 0.05, y - 0.23, abs(x2 - x1) - 0.1, 0.2, label, 11, False, MUTED, PP_ALIGN.CENTER)
        y += 0.48
    footer(s, 8)


def slide_activity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "UML - Activité réservation", "La réservation suit un workflow contrôlé jusqu'à la facture.")
    steps = [
        ("Choisir annonce", 0.95, 2.0),
        ("Demander réservation", 3.0, 2.0),
        ("Notification vendeur", 5.25, 2.0),
        ("Approuver ou rejeter", 7.55, 2.0),
        ("Livraison / remise", 9.85, 1.45),
        ("Commande finalisée", 9.85, 2.75),
        ("Facture PDF", 7.55, 4.0),
    ]
    for label, x, y in steps:
        node(s, x, y, 1.75, 0.58, label, WHITE, PURPLE, 13)
    for a, b in [((2.7, 2.29), (3.0, 2.29)), ((4.75, 2.29), (5.25, 2.29)), ((7.0, 2.29), (7.55, 2.29)), ((9.3, 2.12), (9.85, 1.75)), ((9.3, 2.45), (9.85, 3.05)), ((9.85, 3.34), (8.45, 4.0))]:
        line(s, a[0], a[1], b[0], b[1], PURPLE, 1.5, True)
    text_box(s, 7.65, 2.77, 1.6, 0.24, "rejet: annonce rouverte", 11, False, ORANGE, PP_ALIGN.CENTER)
    footer(s, 9)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Architecture logicielle", "YaZoo sépare clairement interface, API, persistance et services temps réel.")
    nodes = [
        ("React 19 + Vite", 0.75, 2.0, PURPLE),
        ("Axios + Router", 0.75, 3.05, PURPLE),
        ("Laravel 12 API", 4.1, 2.0, BLUE),
        ("Sanctum Auth", 4.1, 3.05, BLUE),
        ("MySQL", 7.3, 1.65, GREEN),
        ("Redis", 7.3, 2.65, GREEN),
        ("Storage médias", 7.3, 3.65, GREEN),
        ("Notifications", 10.15, 2.65, ORANGE),
    ]
    for label, x, y, c in nodes:
        node(s, x, y, 2.1, 0.58, label, WHITE, c, 14)
    for a, b in [((2.85,2.29),(4.1,2.29)), ((2.85,3.34),(4.1,3.34)), ((6.2,2.29),(7.3,1.94)), ((6.2,2.54),(7.3,2.94)), ((6.2,2.84),(7.3,3.94)), ((9.4,2.94),(10.15,2.94))]:
        line(s, a[0], a[1], b[0], b[1], LINE, 1.5, True)
    card(s, 0.75, 5.0, 11.55, 0.8, LILAC, LINE)
    text_box(s, 1.05, 5.18, 10.9, 0.36, "Principe: une API REST sécurisée, consommée par une SPA React, avec cache/queues Redis et base relationnelle MySQL.", 17, True, PURPLE)
    footer(s, 10)


def slide_database(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Base de données", "Le schéma relationnel couvre les interactions sociales et commerciales.")
    tables = [
        ("users", "posts, stories, animals, products"),
        ("posts", "comments, likes, media"),
        ("animals", "reservations, adoption flag"),
        ("products", "stock, reservations"),
        ("communities", "community_members"),
        ("conversations", "messages"),
        ("reservations", "invoice, review, status"),
        ("notifications", "read_at, payload"),
    ]
    for i, (h, b) in enumerate(tables):
        x = 0.75 + (i % 4) * 3.05
        y = 1.55 + (i // 4) * 1.65
        card(s, x, y, 2.65, 1.0, WHITE, LINE, MSO_SHAPE.RECTANGLE)
        text_box(s, x + 0.12, y + 0.12, 2.35, 0.24, h, 15, True, PURPLE, PP_ALIGN.CENTER)
        text_box(s, x + 0.12, y + 0.45, 2.35, 0.34, b, 12, False, MUTED, PP_ALIGN.CENTER)
    text_box(s, 0.9, 5.25, 11.0, 0.35, "Index ajoutés: recherche marketplace, recherche communauté, colonnes de scaling pour améliorer les performances.", 17, True, INK)
    footer(s, 11)


def slide_ui_showcase(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Réalisation UI", "L'interface YaZoo adopte une identité douce, claire et responsive.")
    picture(s, SHOTS[0], 0.8, 1.45, 5.4, 3.05)
    picture(s, SHOTS[11], 6.55, 1.45, 2.05, 3.05)
    picture(s, SHOTS[4], 8.95, 1.45, 3.45, 3.05)
    bullet_list(s, 0.95, 5.0, 11.3, 0.9, [
        "Design violet, composants réutilisables, navigation desktop et mobile.",
        "Pages publiques, authentification, feed, profil, marketplace, communautés et historique.",
    ], 16)
    footer(s, 12)


def slide_auth(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Authentification", "Sanctum sécurise les accès et garde une expérience fluide.")
    picture(s, SHOTS[2], 0.8, 1.45, 5.55, 3.1)
    picture(s, SHOTS[3], 6.75, 1.45, 5.55, 3.1)
    bullet_list(s, 1.0, 4.95, 10.8, 0.95, [
        "Inscription avec validation des champs et création de profil.",
        "Connexion via API Laravel, cookie sécurisé et récupération du profil courant.",
    ], 16)
    footer(s, 13)


def slide_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Fonctionnalités principales", "Le coeur produit combine feed social, profil et marketplace.")
    picture(s, SHOTS[4], 0.65, 1.45, 3.85, 2.15)
    picture(s, SHOTS[5], 4.75, 1.45, 3.85, 2.15)
    picture(s, SHOTS[9], 8.85, 1.45, 3.85, 2.15)
    labels = [
        ("Feed social", "Posts, médias, stories, likes, commentaires."),
        ("Profil", "Bio, statistiques, publications récentes."),
        ("Marketplace", "Annonces animaux/produits, filtres, réservations."),
    ]
    for i, (h, b) in enumerate(labels):
        x = 0.65 + i * 4.1
        text_box(s, x, 3.85, 3.65, 0.28, h, 18, True, PURPLE)
        text_box(s, x, 4.22, 3.65, 0.55, b, 16, False, MUTED)
    footer(s, 14)


def slide_mobile(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Responsive design", "Les parcours restent utilisables sur mobile grâce à une navigation adaptée.")
    picture(s, SHOTS[8], 1.15, 1.35, 2.55, 4.55)
    picture(s, SHOTS[10], 4.15, 1.35, 2.55, 4.55)
    picture(s, SHOTS[9], 7.15, 1.35, 2.55, 4.55)
    text_box(s, 10.15, 1.85, 2.2, 1.2, "Menu latéral, barre basse, cartes verticales et formulaires adaptés aux petits écrans.", 17, True, INK)
    footer(s, 15)


def slide_api(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "API Laravel", "Les routes REST structurent chaque domaine métier.")
    groups = [
        ("Auth", "/auth/register\n/auth/login\n/auth/me"),
        ("Feed", "/posts\n/stories\n/comments"),
        ("Marketplace", "/animals\n/products\n/reservations"),
        ("Communautés", "/communities\nmembership requests"),
        ("Messagerie", "/conversations\n/messages"),
        ("Admin", "/admin/orders\n/admin/moderation"),
    ]
    for i, (h, b) in enumerate(groups):
        x = 0.8 + (i % 3) * 4.15
        y = 1.55 + (i // 3) * 1.75
        card(s, x, y, 3.65, 1.18, WHITE, LINE)
        text_box(s, x + 0.2, y + 0.14, 3.2, 0.25, h, 17, True, PURPLE)
        text_box(s, x + 0.2, y + 0.48, 3.2, 0.42, b, 13, False, MUTED)
    footer(s, 16)


def slide_security(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Sécurité", "Le projet intègre des protections applicatives et cloud.")
    bullet_list(s, 0.9, 1.55, 5.6, 3.6, [
        "Authentification Laravel Sanctum et middleware auth:sanctum.",
        "CORS limité à l'URL frontend de production.",
        "Cookies sécurisés, SameSite none en production HTTPS.",
        "Headers de sécurité: CSP, X-Frame-Options, Referrer-Policy.",
        "Secrets exclus du dépôt et stockés dans Azure Key Vault.",
    ], 16)
    card(s, 7.0, 1.7, 4.75, 2.65, LILAC, LINE)
    text_box(s, 7.35, 2.05, 4.05, 0.42, "Contrôles qualité", 22, True, PURPLE)
    bullet_list(s, 7.35, 2.65, 3.95, 1.3, ["PHPUnit", "Vitest", "ESLint", "Composer audit", "npm audit", "SonarQube"], 16)
    footer(s, 17)


def slide_docker(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Docker", "L'environnement regroupe les services nécessaires au développement et à la production.")
    comps = [
        ("frontend", "Nginx + build React", 0.9, 1.7),
        ("nginx", "Reverse proxy API", 3.8, 1.7),
        ("app", "Laravel PHP-FPM", 6.7, 1.7),
        ("queue", "Worker Redis", 9.6, 1.7),
        ("mysql", "MySQL 8.4", 3.0, 3.75),
        ("redis", "Redis cache/queue", 6.1, 3.75),
        ("sonarqube", "Qualité code", 9.2, 3.75),
    ]
    for name, desc, x, y in comps:
        node(s, x, y, 2.25, 0.68, name, WHITE, PURPLE, 15)
        text_box(s, x, y + 0.78, 2.25, 0.28, desc, 12, False, MUTED, PP_ALIGN.CENTER)
    for a, b in [((3.15,2.04),(3.8,2.04)), ((6.05,2.04),(6.7,2.04)), ((8.95,2.04),(9.6,2.04)), ((7.8,2.38),(4.1,3.75)), ((7.8,2.38),(7.2,3.75))]:
        line(s, a[0], a[1], b[0], b[1], LINE, 1.5, True)
    footer(s, 18)


def slide_ci_cd(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "CI/CD", "Le pipeline automatise les contrôles avant le déploiement.")
    steps = [
        ("Push GitHub", "main"),
        ("Tests backend", "PHPUnit"),
        ("Tests frontend", "Vitest + ESLint"),
        ("Build Docker", "API / Front"),
        ("Registry", "Docker Hub"),
        ("Azure", "App Service"),
    ]
    x = 0.75
    for i, (h, b) in enumerate(steps):
        node(s, x, 2.45, 1.65, 0.72, h, WHITE, PURPLE, 13)
        text_box(s, x, 3.28, 1.65, 0.25, b, 12, False, MUTED, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            line(s, x + 1.65, 2.81, x + 2.05, 2.81, PURPLE, 1.5, True)
        x += 2.05
    card(s, 1.2, 4.65, 10.9, 0.75, LILAC, LINE)
    text_box(s, 1.55, 4.82, 10.2, 0.32, "Résultat: moins d'erreurs manuelles, builds reproductibles et déploiement cloud traçable.", 17, True, PURPLE)
    footer(s, 19)


def slide_azure(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Déploiement Azure", "La version cloud tourne sur App Service avec MySQL, Redis et Key Vault.")
    items = [
        ("Frontend", "yazoo.azurewebsites.net", 0.8, 1.55, PURPLE),
        ("Backend API", "yazoo-api.azurewebsites.net", 4.15, 1.55, BLUE),
        ("MySQL", "yazoo-mysql-0c2b09", 7.5, 1.55, GREEN),
        ("Managed Redis", "yazoo-redis-0c2b09", 7.5, 3.0, GREEN),
        ("Key Vault", "yazoo-kv-0c2b09", 4.15, 4.45, ORANGE),
        ("Docker Hub", "5eef/yazoo-*", 0.8, 4.45, ORANGE),
    ]
    for h, b, x, y, c in items:
        node(s, x, y, 2.75, 0.65, h, WHITE, c, 15)
        text_box(s, x, y + 0.74, 2.75, 0.3, b, 11, False, MUTED, PP_ALIGN.CENTER)
    for a, b in [((3.55,1.88),(4.15,1.88)), ((6.9,1.88),(7.5,1.88)), ((6.9,2.08),(7.5,3.33)), ((2.2,4.45),(4.15,2.1)), ((5.52,4.45),(5.52,2.2))]:
        line(s, a[0], a[1], b[0], b[1], LINE, 1.4, True)
    footer(s, 20)


def slide_tests(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Tests et validation", "La préparation production a été validée par des tests automatisés.")
    metrics = [
        ("69", "tests Laravel passés"),
        ("21", "tests React/Vitest passés"),
        ("0", "vulnérabilité npm audit"),
        ("OK", "Docker build API + Front"),
        ("UP", "SonarQube local"),
        ("200", "/health/ready Azure"),
    ]
    for i, (v, l) in enumerate(metrics):
        x = 0.8 + (i % 3) * 4.05
        y = 1.65 + (i // 3) * 1.75
        card(s, x, y, 3.5, 1.1, WHITE, LINE)
        text_box(s, x + 0.2, y + 0.15, 1.1, 0.38, v, 24, True, PURPLE)
        text_box(s, x + 1.15, y + 0.25, 2.0, 0.35, l, 16, False, INK)
    footer(s, 21)


def slide_gantt(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Planification", "Le projet a été conduit par étapes, de l'analyse au déploiement.")
    phases = [
        ("Analyse", 1.0, 1.75, 2.0, PURPLE),
        ("Conception UML", 2.4, 2.35, 2.2, BLUE),
        ("Backend Laravel", 3.6, 2.95, 3.0, GREEN),
        ("Frontend React", 4.5, 3.55, 3.0, PURPLE),
        ("Tests", 7.0, 4.15, 2.0, ORANGE),
        ("Docker/Azure", 8.25, 4.75, 2.65, BLUE),
    ]
    for month, x in [("S1", 1.0), ("S2", 2.8), ("S3", 4.6), ("S4", 6.4), ("S5", 8.2), ("S6", 10.0)]:
        text_box(s, x, 1.25, 0.5, 0.25, month, 13, True, MUTED, PP_ALIGN.CENTER)
        line(s, x + 0.25, 1.55, x + 0.25, 5.6, LINE, 0.8)
    for label, x, y, w, c in phases:
        pill(s, x, y, w, 0.42, label, c)
    footer(s, 22)


def slide_bilan(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Bilan", "YaZoo répond au besoin initial avec une plateforme complète et déployée.")
    cols = [
        ("Acquis techniques", ["Laravel API", "React SPA", "Docker", "CI/CD", "Azure Cloud"]),
        ("Apports métier", ["Communauté", "Adoption", "Marketplace", "Messagerie", "Réservations"]),
        ("Difficultés traitées", ["CORS/Sanctum", "Uploads médias", "Queues Redis", "Policies Azure", "Déploiement Docker"]),
    ]
    for i, (h, items) in enumerate(cols):
        x = 0.8 + i * 4.15
        card(s, x, 1.55, 3.7, 3.65, WHITE, LINE)
        text_box(s, x + 0.22, 1.85, 3.2, 0.32, h, 18, True, PURPLE)
        bullet_list(s, x + 0.28, 2.42, 3.0, 2.2, items, 16)
    footer(s, 23)


def slide_perspectives(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title(s, "Perspectives", "Les prochaines évolutions peuvent renforcer la valeur de YaZoo.")
    bullet_list(s, 0.9, 1.55, 5.6, 3.5, [
        "Application mobile ou PWA avancée.",
        "Paiement en ligne et suivi de livraison.",
        "Géolocalisation des annonces et refuges.",
        "Tableau analytique admin plus détaillé.",
        "Messagerie temps réel activée via WebSockets.",
        "Stockage cloud dédié pour les médias.",
    ], 17)
    picture(s, SHOTS[7], 7.0, 1.55, 4.95, 2.85)
    card(s, 7.0, 4.75, 4.95, 0.75, LILAC, LINE)
    text_box(s, 7.35, 4.92, 4.25, 0.34, "Objectif futur: passer d'une plateforme fonctionnelle à un écosystème animalier complet.", 16, True, PURPLE)
    footer(s, 24)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    picture(s, LOGO_YAZOO, 5.65, 0.85, 1.75, 1.75)
    text_box(s, 2.45, 2.8, 8.5, 0.55, "Merci pour votre attention", 28, True, PURPLE, PP_ALIGN.CENTER)
    text_box(s, 2.2, 3.55, 9.0, 0.45, "Questions / Discussion", 22, True, INK, PP_ALIGN.CENTER)
    text_box(s, 2.45, 4.35, 8.5, 0.55, "Youssef Boughioul | Salma Kabran | Chams Doha Amine", 17, False, MUTED, PP_ALIGN.CENTER)
    text_box(s, 2.45, 4.82, 8.5, 0.35, "Encadrante: Imane LAASSAR", 17, False, MUTED, PP_ALIGN.CENTER)
    footer(s, 25)


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_agenda(prs)
    slide_context(prs)
    slide_objectives(prs)
    slide_stack(prs)
    slide_use_case(prs)
    slide_class(prs)
    slide_sequence(prs)
    slide_activity(prs)
    slide_architecture(prs)
    slide_database(prs)
    slide_ui_showcase(prs)
    slide_auth(prs)
    slide_features(prs)
    slide_mobile(prs)
    slide_api(prs)
    slide_security(prs)
    slide_docker(prs)
    slide_ci_cd(prs)
    slide_azure(prs)
    slide_tests(prs)
    slide_gantt(prs)
    slide_bilan(prs)
    slide_perspectives(prs)
    slide_thanks(prs)
    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
